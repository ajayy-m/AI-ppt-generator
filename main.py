import os
import json
import random
import requests
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from google.generativeai import configure, GenerativeModel
from PIL import Image

# === CONFIGURATION ===
GOOGLE_API_KEY = "Enter_your_Api"  # Replace this with your actual Gemini API key
SEARCH_ENGINE_ID = "Google_search_engine API"  # Replace with your Google Custom Search Engine ID
GEMINI_API = "Gemini API key - Register/get it, and then add it here"

configure(api_key=GEMINI_API)
model = GenerativeModel(model_name="models/gemini-1.5-pro")

TEMPLATE_DIR = "templates"

# === UTILS ===
def extract_valid_json(text):
    match = re.search(r'\[.*\]', text, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            return []
    return []

def generate_subtopics(topic, count=5, retries=2):
    assert count >= 3, "Slide count must be at least 3 to include intro and conclusion."
    middle = count - 2
    prompt = f'''
You are an expert presentation assistant.

Generate a JSON list of {count} slides on the topic: "{topic}", structured as follows:
1. "Introduction" — a slide introducing the topic
2. {middle} subtopic slides
3. "Conclusion" — a summary slide

Each slide must include:
- "title"
- "content": either a paragraph or a list of 3-5 bullet points
- "image_type": one of ["diagram", "illustration", "chart", "photo", "concept"]

Only return valid JSON. No commentary.
'''
    for attempt in range(retries):
        response = model.generate_content(prompt)
        slides = extract_valid_json(response.text)
        if slides:
            return slides
        print(f"⚠️ Attempt {attempt + 1}: Failed to get valid JSON from Gemini.")
    print("❌ Gemini failed to return valid slide data after retries.")
    return []

def fetch_image(main_topic, subtopic, image_type="concept"):
    query = f"{main_topic} {subtopic} {image_type} high quality realistic professional"
    search_url = f"https://www.googleapis.com/customsearch/v1?q={query}&searchType=image&key={GOOGLE_API_KEY}&cx={SEARCH_ENGINE_ID}&num=3"
    try:
        response = requests.get(search_url).json()
        items = response.get("items", [])
        for item in items:
            img_url = item.get("link")
            try:
                img_response = requests.get(img_url, timeout=5)
                content_type = img_response.headers.get("Content-Type", "")
                if "image" not in content_type:
                    continue
                img_data = img_response.content
                img = Image.open(io.BytesIO(img_data))
                if img.format == "WEBP":
                    img = img.convert("RGB")
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    return buf
                return io.BytesIO(img_data)
            except Exception as e:
                print(f"⚠️ Skipping bad image URL: {img_url} — {e}")
    except Exception as e:
        print(f"⚠️ Failed image search: {e}")
    return None

def create_content_slide(prs, title, content, image_stream, index):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    margin = Inches(0.5)
    image_width = Inches(4.5)
    image_height = Inches(4.0)
    text_width = Inches(4.5)
    text_height = Inches(5.0)

    if index % 2 == 0:
        image_left = margin
        text_left = Inches(5)
    else:
        image_left = Inches(5)
        text_left = margin

    if image_stream:
        try:
            slide.shapes.add_picture(image_stream, image_left, Inches(1.0), width=image_width, height=image_height)
        except Exception as e:
            print(f"⚠️ Could not add image: {e}")

    textbox = slide.shapes.add_textbox(text_left, Inches(0.5), text_width, text_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = True

    # Title
    title_p = tf.paragraphs[0]
    title_p.text = title
    title_p.font.bold = True
    title_p.font.size = Pt(28)
    title_p.alignment = PP_ALIGN.LEFT

    # Content
    if isinstance(content, list):
        for bullet in content:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1
            p.font.size = Pt(18)
    else:
        p = tf.add_paragraph()
        p.text = content
        p.level = 1
        p.font.size = Pt(18)

def add_title_slide(prs, topic):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(2))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def choose_template():
    templates = [f for f in os.listdir(TEMPLATE_DIR) if f.endswith(".pptx")]
    if not templates:
        raise FileNotFoundError("No PPTX templates found in 'templates/' folder.")
    print("Available templates:")
    for idx, t in enumerate(templates, 1):
        print(f"{idx}. {t}")
    choice = int(input("Choose a template by number: ")) - 1
    return os.path.join(TEMPLATE_DIR, templates[choice])

# === MAIN ===
def main():
    topic = input("Enter your presentation topic: ")
    slide_count = int(input("How many slides (min 3)? "))

    slides = generate_subtopics(topic, slide_count)
    if not slides:
        return

    # Reorder with fallback if intro/conclusion missing
    intro = next((s for s in slides if s["title"].strip().lower() == "introduction"), None)
    conclusion = next((s for s in slides if s["title"].strip().lower() == "conclusion"), None)

    if not intro:
        intro = slides[0]

    if not conclusion:
        conclusion = slides[-1]

    middle = [s for s in slides if s not in [intro, conclusion]]
    ordered_slides = [intro] + middle + [conclusion]

    template_path = choose_template()
    prs = Presentation(template_path)

    add_title_slide(prs, topic)

    for idx, slide in enumerate(ordered_slides):
        title = slide["title"]
        content = slide["content"]
        image_type = slide.get("image_type", "concept")
        img = fetch_image(topic, title, image_type)
        create_content_slide(prs, title, content, img, idx)

    output_path = f"{topic.replace(' ', '_')}.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}")

if __name__ == "__main__":
    main()
 
